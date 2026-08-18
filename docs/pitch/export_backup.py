"""Static backups of the HTML deck: docs/VayuNetra_Pitch.pptx (one full-bleed picture per
slide + the speaker notes) and docs/VayuNetra_Pitch.pdf. Needs the renders from
`node web/scripts/qa/deck-render.mjs` (static mode) in docs/pitch/render/.

    python3 docs/pitch/export_backup.py        # system python3 has python-pptx + Pillow
"""
from __future__ import annotations

import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

HERE = Path(__file__).resolve().parent
ROOT = HERE.parents[1]
RENDER = HERE / "render"
DECK = HERE / "deck.html"


def notes() -> list[str]:
    html = DECK.read_text()
    out = []
    for m in re.finditer(r'<section class="slide"([^>]*)>', html):
        attrs = m.group(1)
        n = re.search(r'data-notes="([^"]*)"', attrs)
        out.append((n.group(1) if n else "").replace("&gt;", ">").replace("&amp;", "&"))
    return out


def main() -> None:
    files = sorted(RENDER.glob("slide-*.png"))
    if not files:
        raise SystemExit("no renders — run: cd web && node scripts/qa/deck-render.mjs")
    nts = notes()
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    jpg_dir = RENDER / "jpg"
    jpg_dir.mkdir(exist_ok=True)
    for i, f in enumerate(files):
        j = jpg_dir / (f.stem + ".jpg")
        Image.open(f).convert("RGB").save(j, quality=86, optimize=True)
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(j), 0, 0, width=prs.slide_width, height=prs.slide_height)
        s.notes_slide.notes_text_frame.text = nts[i] if i < len(nts) else ""
    out_pptx = ROOT / "docs" / "VayuNetra_Pitch.pptx"
    prs.save(out_pptx)
    imgs = [Image.open(f).convert("RGB") for f in files]
    out_pdf = ROOT / "docs" / "VayuNetra_Pitch.pdf"
    imgs[0].save(out_pdf, save_all=True, append_images=imgs[1:], resolution=144)
    print(f"wrote {out_pptx} ({out_pptx.stat().st_size/1e6:.1f} MB) and {out_pdf} ({out_pdf.stat().st_size/1e6:.1f} MB), {len(files)} slides")


if __name__ == "__main__":
    main()
